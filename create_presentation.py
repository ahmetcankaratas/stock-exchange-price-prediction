from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Turkish Airlines Stock Price Prediction"
    subtitle.text = "Following CRISP-DM Methodology"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Format subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add content
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.level = 0

def main():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Create title slide
    create_title_slide(prs)
    
    # Business Understanding
    business_understanding = [
        "Project Objectives:",
        "• Predict Turkish Airlines (THY) stock prices using machine learning",
        "• Develop a reliable forecasting model for investment decision support",
        "• Understand key factors influencing THY stock performance",
        "",
        "Success Criteria:",
        "• Model accuracy in predicting stock price movements",
        "• Ability to capture market trends and patterns",
        "• Practical usability for investment decisions"
    ]
    create_content_slide(prs, "1. Business Understanding", business_understanding)
    
    # Data Understanding
    data_understanding = [
        "Data Sources:",
        "• Yahoo Finance API via RapidAPI",
        "• Historical stock data at multiple intervals (5m, 1h, 1d)",
        "• General market data for context",
        "",
        "Data Types:",
        "• Stock prices (Open, High, Low, Close)",
        "• Trading volume",
        "• Market indicators",
        "• Technical indicators",
        "",
        "Data Quality Assessment:",
        "• Missing value analysis",
        "• Outlier detection",
        "• Data consistency checks",
        "• Time series completeness"
    ]
    create_content_slide(prs, "2. Data Understanding", data_understanding)
    
    # Data Preparation
    data_preparation = [
        "Data Collection:",
        "• Automated data fetching through RapidAPI",
        "• Multiple time intervals for comprehensive analysis",
        "• Regular data updates",
        "",
        "Feature Engineering:",
        "• Moving Averages (5-day and 20-day)",
        "• Relative Strength Index (RSI)",
        "• MACD (Moving Average Convergence Divergence)",
        "• Bollinger Bands",
        "• Volume-based indicators",
        "",
        "Data Preprocessing:",
        "• Data cleaning and normalization",
        "• Handling missing values",
        "• Time series alignment",
        "• Feature scaling"
    ]
    create_content_slide(prs, "3. Data Preparation", data_preparation)
    
    # Modeling
    modeling = [
        "Model Architecture:",
        "• LSTM (Long Short-Term Memory) neural network",
        "• Multiple LSTM layers with dropout",
        "• 60-day lookback period",
        "• Early stopping implementation",
        "",
        "Model Training:",
        "• Training/validation/test split",
        "• Hyperparameter optimization",
        "• Cross-validation",
        "• Model persistence",
        "",
        "Technical Implementation:",
        "• Python-based implementation",
        "• TensorFlow/Keras framework",
        "• GPU acceleration support",
        "• Modular code structure"
    ]
    create_content_slide(prs, "4. Modeling", modeling)
    
    # Evaluation
    evaluation = [
        "Model Performance Metrics:",
        "• Mean Absolute Error (MAE)",
        "• Root Mean Square Error (RMSE)",
        "• Directional Accuracy",
        "• R-squared Score",
        "",
        "Validation Methods:",
        "• Backtesting on historical data",
        "• Cross-validation results",
        "• Performance on different market conditions",
        "• Comparison with baseline models",
        "",
        "Model Limitations:",
        "• Market unpredictability",
        "• External factors impact",
        "• API rate limits",
        "• Computational requirements"
    ]
    create_content_slide(prs, "5. Evaluation", evaluation)
    
    # Deployment
    deployment = [
        "Implementation:",
        "• Automated data collection pipeline",
        "• Regular model updates",
        "• Prediction API endpoints",
        "• Monitoring system",
        "",
        "Maintenance:",
        "• Regular data updates",
        "• Model retraining schedule",
        "• Performance monitoring",
        "• Error handling",
        "",
        "Future Improvements:",
        "• Additional feature engineering",
        "• Alternative model architectures",
        "• Real-time prediction capabilities",
        "• Enhanced visualization tools"
    ]
    create_content_slide(prs, "6. Deployment", deployment)
    
    # Project Structure
    project_structure = [
        "Project Structure:",
        "project/",
        "├── data/",
        "│   ├── stocks.csv",
        "│   └── thy_historical_{interval}.csv",
        "├── notebooks/",
        "├── src/",
        "│   ├── data_collection.py",
        "│   ├── data_preprocessing.py",
        "│   ├── model.py",
        "│   ├── utils.py",
        "│   └── visualization.py",
        "└── requirements.txt"
    ]
    create_content_slide(prs, "Project Structure", project_structure)
    
    # Next Steps
    next_steps = [
        "Next Steps:",
        "1. Implement additional technical indicators",
        "2. Explore ensemble methods",
        "3. Develop real-time prediction capabilities",
        "4. Create interactive visualization dashboard",
        "5. Optimize model performance"
    ]
    create_content_slide(prs, "Next Steps", next_steps)
    
    # Contact
    contact = [
        "Contact:",
        "For questions and collaboration:",
        "• GitHub Repository: stock-exchange-price-prediction",
        "• Project Documentation: README.md"
    ]
    create_content_slide(prs, "Contact", contact)
    
    # Save the presentation
    prs.save('Turkish_Airlines_Stock_Prediction.pptx')

if __name__ == '__main__':
    main() 